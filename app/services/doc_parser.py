"""VLM 기반 문서 파싱 모듈.

지원 형식: .pptx/.ppt, .docx/.doc, .xlsx/.xls, .pdf, .md/.txt

처리 흐름:
  파일 바이트  →  형식별 파서  →  텍스트/이미지 추출
                                    ↓
                      이미지 → Ollama VLM(llava) → 설명 텍스트
                                    ↓
                              청크(list[str]) 반환
"""
from __future__ import annotations
import base64
import io
from pathlib import Path
from typing import Any

from app.lib.ollama import OllamaClient
from app.config import settings

CHUNK_SIZE    = 800
CHUNK_OVERLAP = 100

# ── 공통 유틸 ─────────────────────────────────────────────────────────────────

def _chunk_text(text: str) -> list[str]:
    words = text.split()
    chunks: list[str] = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i : i + CHUNK_SIZE])
        if chunk.strip():
            chunks.append(chunk)
        i += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks


async def _vlm_describe(
    image_bytes: bytes,
    ollama: OllamaClient,
    hint: str = "이 이미지의 내용을 한국어로 상세히 설명해줘.",
) -> str:
    """
    Ollama VLM(설정: VLM_MODEL)을 이용해 이미지를 한국어로 설명한다.
    VLM 모델이 없거나 오류 시 빈 문자열을 반환한다.
    """
    try:
        encoded = base64.b64encode(image_bytes).decode()
        result = await ollama.chat(
            settings.VLM_MODEL,
            [{"role": "user", "content": hint, "images": [encoded]}],
            {"temperature": 0.1, "num_predict": 512},
        )
        return result.strip()
    except Exception:
        return ""


# ── 형식별 파서 ───────────────────────────────────────────────────────────────

async def parse_pptx(file_bytes: bytes, ollama: OllamaClient) -> list[str]:
    """PowerPoint(.pptx/.ppt): 슬라이드 텍스트 + 이미지 VLM 설명."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs    = Presentation(io.BytesIO(file_bytes))
    chunks: list[str] = []

    for idx, slide in enumerate(prs.slides, 1):
        parts: list[str] = []

        for shape in slide.shapes:
            # 텍스트 프레임
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)

            # 이미지 → VLM
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    desc = await _vlm_describe(shape.image.blob, ollama)
                    if desc:
                        parts.append(f"[이미지] {desc}")
                except Exception:
                    pass

        if parts:
            content = f"[슬라이드 {idx}]\n" + "\n".join(parts)
            chunks.extend(_chunk_text(content))

    return chunks


async def parse_docx(file_bytes: bytes, ollama: OllamaClient) -> list[str]:
    """Word(.docx/.doc): 단락 + 표 + 이미지 VLM 설명."""
    from docx import Document as DocxDocument

    doc   = DocxDocument(io.BytesIO(file_bytes))
    parts: list[str] = []

    # 단락
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # 표
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append("[표]\n" + "\n".join(rows))

    # 이미지 관계(rels)에서 추출
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                desc = await _vlm_describe(rel.target_part.blob, ollama)
                if desc:
                    parts.append(f"[이미지] {desc}")
            except Exception:
                pass

    full_text = "\n\n".join(parts)
    return _chunk_text(full_text) if full_text.strip() else []


async def parse_xlsx(file_bytes: bytes, ollama: OllamaClient) -> list[str]:
    """Excel(.xlsx/.xls): 시트별 셀 값을 텍스트로 변환."""
    import openpyxl

    wb     = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    chunks: list[str] = []

    for sheet_name in wb.sheetnames:
        ws    = wb[sheet_name]
        rows  = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            content = f"[시트: {sheet_name}]\n" + "\n".join(rows)
            chunks.extend(_chunk_text(content))

    return chunks


async def parse_pdf(file_bytes: bytes, ollama: OllamaClient) -> list[str]:
    """PDF: 텍스트 우선 추출, 텍스트 희박 페이지는 VLM으로 설명."""
    import pdfplumber

    chunks: list[str] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for idx, page in enumerate(pdf.pages, 1):
            text = (page.extract_text() or "").strip()

            if len(text) < 50:
                # 이미지 중심 페이지: VLM으로 처리
                try:
                    img_buf = io.BytesIO()
                    page.to_image(resolution=150).save(img_buf, format="PNG")
                    text = await _vlm_describe(
                        img_buf.getvalue(), ollama,
                        f"PDF {idx}페이지 내용을 한국어로 설명해줘.",
                    )
                except Exception:
                    pass

            if text.strip():
                content = f"[페이지 {idx}]\n{text}"
                chunks.extend(_chunk_text(content))

    return chunks


async def parse_text(file_bytes: bytes, ollama: OllamaClient) -> list[str]:
    """Plain text / Markdown: 그대로 청킹."""
    text = file_bytes.decode("utf-8", errors="replace")
    return _chunk_text(text)


# ── 라우터 ────────────────────────────────────────────────────────────────────

_PARSERS: dict[str, Any] = {
    ".pptx": parse_pptx,
    ".ppt":  parse_pptx,
    ".docx": parse_docx,
    ".doc":  parse_docx,
    ".xlsx": parse_xlsx,
    ".xls":  parse_xlsx,
    ".pdf":  parse_pdf,
    ".md":   parse_text,
    ".txt":  parse_text,
}

SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(_PARSERS.keys())


async def parse_document(
    filename:   str,
    file_bytes: bytes,
    ollama:     OllamaClient,
) -> list[str]:
    """
    파일명의 확장자를 보고 적절한 파서를 선택하여 청크 목록을 반환한다.

    Raises:
        ValueError: 지원하지 않는 파일 형식
    """
    ext    = Path(filename).suffix.lower()
    parser = _PARSERS.get(ext)
    if not parser:
        raise ValueError(
            f"지원하지 않는 파일 형식: {ext}  "
            f"(지원: {', '.join(sorted(SUPPORTED_EXTENSIONS))})"
        )
    return await parser(file_bytes, ollama)
